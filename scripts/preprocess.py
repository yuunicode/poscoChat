# scripts/preprocess.py
import logging, json, uuid, shutil, os
from pathlib import Path
from config.enums import ChunkingType, ChunkingStrategy
from config.settings import get_settings
from infra.chunker import Chunker
from abc import ABC, abstractmethod
from typing import List, Dict, Any, Optional
from unstructured.partition.docx import partition_docx
from docx import Document  # DOCX 파일 처리
from pptx import Presentation  # PPTX 파일 처리
from pptx.table import Table
from pptx.enum.shapes import MSO_SHAPE_TYPE

# settings 객체 가져오기, 로깅 설정
settings = get_settings() 
logging.basicConfig(level = logging.INFO, format = '%(asctime)s - %(levelname)s - %(message)s')


class BasePreprocessor(ABC):
    def __init__(self, filename: str):
        """ BasePreprocessor는 모든 전처리기 클래스의 기본 클래스입니다. """
        self.logger = logging.getLogger(__name__)
        self.filepath = settings.RAW_DATA_PATH / filename
        self.records = [] # 초기화 시엔 비어있으며 parse 메서드 콜 이후로 채워짐
        self.chunker = Chunker()
            
    @abstractmethod
    def parse(self) -> List[Dict[str, Any]]:
        """
        파일을 파싱하여 레코드를 생성하는 메소드.
        
        Args:
            file_path (Path): 전처리할 파일의 경로.
        
        Returns:
            List[Dict[str, Any]]: 전처리된 레코드 목록.
        """
        pass

    @abstractmethod
    def preprocess(self) -> List[Dict[str, Any]]:
        """
        전처리 작업을 정의하는 추상 메서드입니다.
        
        Args:
            file_path (Path): 전처리할 파일의 경로.
        
        Returns:
            List[Dict[str, Any]]: 전처리된 레코드 목록.
        """
        pass

    @abstractmethod
    def get_chunking_configs(self) -> List[Dict[str, Any]]:
        """
        이 전처리기 인스턴스에 특화된 청킹 설정 목록을 반환합니다.
        각 설정은 'chunking_type', 'strategy', 그리고 선택적인 'kwargs'를 포함해야 합니다.
        """
        pass

    @abstractmethod
    def postprocess(self) -> List[Dict[str, Any]]:
        """
        후처리 작업을 정의하는 추상 메서드입니다.
        
        Args:
            file_path (Path): 후처리할 파일의 경로.
        
        Returns:
            None
        """
        pass

    def chunk(self, chunking_type: ChunkingType, strategy: ChunkingStrategy, **kwargs) -> None:
        
        text_records = [rec for rec in self.records if rec.get("category") != "Table"]
        table_records = [rec for rec in self.records if rec.get("category") == "Table"]
        processed_records: List[Dict[str, Any]] = []
        self.logger.info(f"청킹 전 현재 텍스트 레코드: {len(text_records)}, 테이블 레코드: {len(table_records)}")

        if chunking_type == ChunkingType.NARRATIVE:
            # Apply text chunking only to narrative (non-table) records
            processed_text_records = self.chunker.apply_text_chunking(
                text_records,
                strategy=strategy, # Use the strategy passed to the chunk method
                **kwargs
            )
            self.logger.info(f"텍스트 청킹 완료. 전처리 된 레코드 수: {len(processed_text_records)}")
            processed_records.extend(processed_text_records)
            processed_records.extend(table_records)
            self.logger.info(f'청킹 완료. 총 레코드 수: {len(processed_records)}')

        elif chunking_type == ChunkingType.TABLE:
            # Apply table chunking only to table records
            processed_table_records = self.chunker.apply_table_chunking(
                table_records,
                strategy=strategy, # Use the strategy passed to the chunk method
                **kwargs
            )
            self.logger.info(f"테이블 청킹 완료. 전처리 된 레코드 수: {len(processed_table_records)}")
            processed_records.extend(processed_table_records)
            processed_records.extend(text_records)
            self.logger.info(f'청킹 완료. 총 레코드 수: {len(processed_records)}')
        
        else:
            self.logger.warning(f"알 수 없는 청킹 유형입니다: '{chunking_type}'. 원본 레코드를 반환합니다.")
            return 
        
        self.records =  processed_records

    def to_jsonl(self) -> None:
        """
        레코드를 JSONL 형식으로 저장하는 메서드.
        
        Args:
            records (List[Dict[str, Any]]): 저장할 레코드 목록.
            output_path (Path): 저장할 JSONL 파일 경로.
        """
        try:
            jsonl_output_path = Path(settings.PROCESSED_DATA_FILE_PATH)  # Path 객체 확인
            # 경로가 유효한지 확인
            if not jsonl_output_path.parent.exists():
                raise ValueError(f"디렉터리가 존재하지 않습니다: {jsonl_output_path.parent}")
            
            #파일을 쓰기 모드로 열기
            with jsonl_output_path.open("w", encoding="utf-8") as f:
                for rec in self.records:
                    try:
                        # 레코드 직렬화 (json.dumps에서 직렬화할 수 없는 객체 처리)
                        json_str = json.dumps(rec, ensure_ascii=False)
                        f.write(json_str + "\n")
                    except TypeError as e:
                        self.logger.error(f"레코드 직렬화 오류: {rec} - {e}")
            
            self.logger.info(f"경로: {jsonl_output_path} 에 저장 완료했습니다. : 총 청크 수: {len(self.records)}")
            
        except Exception as e:
            self.logger.error(f"Failed to save JSONL file. : {e}")
            raise

    def _save_table_html(self, html_code: str, table_id: str) -> Path:
        """
        테이블 HTML 코드를 .html 파일로 저장하고, 저장된 파일의 경로를 반환합니다.
        output_dir 내부에 'tables' 폴더를 생성하고 저장합니다.
        """
        try:
            # FILENAME에서 확장자 제거하고, 'tables' 폴더 경로 생성
            filename_without_extension = settings.FILENAME.rsplit('.', 1)[0]  # 확장자 제거
            table_out_path = settings.PROCESSED_DATA_PATH / filename_without_extension / 'tables'
            table_out_path.mkdir(parents=True, exist_ok=True)  # 새로운 폴더 생성

            # HTML 파일 경로 설정
            html_filename = f"table_{table_id}.html"
            html_file_path = table_out_path / html_filename

            # HTML 코드 저장
            html_file_path.write_text(html_code, encoding="utf-8")
            logging.debug(f"Table HTML saved to {html_file_path}")       
            return html_file_path

        except Exception as e:
            logging.warning(f"Error saving table HTML: {e}. Returning default path.")
            return Path(settings.PROCESSED_DATA_PATH) / "tables" / "default_table.html" # 정말 기본 경로 

    def _save_images(self) -> None:
        """
            DOCX/PPTX 파일에서 이미지를 추출하여 지정된 경로에 저장합니다.
            output_dir 내부에 'images' 폴더를 생성하고 저장합니다.
            경고: 현재 구현은 'images' 폴더를 매번 새로 생성하므로, 여러 문서를 처리할 경우 이전 이미지가 삭제될 수 있습니다.
            문서별 고유 폴더에 저장하려면 이 함수 내부 로직을 수정해야 합니다.
        """

        # FILENAME에서 확장자 제거하고, 'images' 폴더 경로 생성
        filename_without_extension = settings.FILENAME.rsplit('.', 1)[0]  # FILENAME에서 확장자 제거
        image_out_path = settings.PROCESSED_DATA_PATH / filename_without_extension / 'images'

        # 기존 폴더가 있으면 삭제 후 새로 생성
        if image_out_path.exists() and image_out_path.is_dir():
            shutil.rmtree(image_out_path)  # 기존 폴더 삭제
        image_out_path.mkdir(parents=True, exist_ok=True)  # 새로운 폴더 생성

        # DOCX 파일 처리
        if settings.FILENAME.endswith(".docx"):
            try:
                doc = Document(str(settings.RAW_DATA_FILE_PATH))
            except FileNotFoundError:
                logging.error(f"Image source file {settings.RAW_DATA_FILE_PATH} not found.")
                return
            except Exception as e:
                logging.error(f"Error opening DOCX for image extraction from {settings.RAW_DATA_FILE_PATH.name}: {e}")
                return
 
            # DOCX에서 이미지 추출
            rels = doc.part.rels
            image_idx = 0

            # 이미지 추출 및 저장
            for rel in rels.values():
                if "image" in rel.reltype:
                    image_idx += 1
                    image_data = rel.target_part.blob
                    ext = Path(rel.target_part.partname).suffix  # 이미지 확장자 추출
                    image_name = f"{filename_without_extension}_image_{image_idx}{ext}"
                    (image_out_path / image_name).write_bytes(image_data)  # 이미지 저장

            logging.info(f"Saved {image_idx} images from DOCX to {image_out_path}")

        # PPTX 파일 처리
        elif settings.FILENAME.endswith(".pptx"):
            try:    
                presentation = Presentation(str(settings.RAW_DATA_FILE_PATH))
            except FileNotFoundError:
                logging.error(f"Image source file {settings.RAW_DATA_FILE_PATH} not found.")
                return
            except Exception as e:
                logging.error(f"Error opening PPTX for image extraction from {settings.RAW_DATA_FILE_PATH.name}: {e}")
                return

            image_idx = 0
            # PPTX에서 이미지 추출
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        ext = image.ext
                        image_idx += 1
                        image_name = f"{filename_without_extension}_image_{image_idx}.{ext}"
                        (image_out_path / image_name).write_bytes(image.blob)

        else:
            logging.error("Unsupported file type. Only DOCX and PPTX files are supported.")
            return

class DocxPreprocessor(BasePreprocessor):
    def __init__(self, filename: str):
        super().__init__(filename)

    # ----------------------------------
    #           상속받은 메소드
    # ----------------------------------    
    def parse(self) -> List[Dict[str, Any]]:
        """
            DOCX 파일을 파싱하여 표준화된 DocumentRecord 리스트를 반환합니다.
            테이블의 경우 원본 html을 table_info 딕셔너리의 'text_as_html' 필드에 포함합니다.

            각 DocumentRecord는 다음 형식을 따릅니다:
            {
                "id": "고유 식별자 (str)",
                "text": "문서 섹션의 텍스트 내용 (str)",
                "metadata": {
                    "general_info": {},  # 초기에는 비어있는 딕셔너리. 이후 단계에서 일반 정보가 추가됨.
                    "table_info": {}     # 초기에는 비어있는 딕셔너리. 이후 단계에서 테이블 정보가 추가됨.
                }
            }
            'metadata' 필드는 이후 전처리 단계에서 'general_info'나 'table_info'와 같은
            추가적인 딕셔너리 형태의 정보를 포함하도록 확장됩니다.

            Args:
                file_path (str): 파싱할 DOCX 파일의 경로.

            Returns:
                Tuple[List[Dict[str, Any]], Dict[str, Any]]:
                    - records: 파싱된 DocumentRecord 객체들의 리스트.
                    - table_info_cache: 테이블 ID를 키로, 테이블의 초기 정보를 값으로 하는 딕셔너리.
        """   
        # ------------- STEP 1: 파일 읽기
        try:
            elements = []
            with open(self.filepath, 'rb') as file:
                elements = partition_docx(file                    = file
                                        , extract_images_in_docx  = False
                                        , include_page_breaks     = False)
        except FileNotFoundError:
            logging.error(f"File {self.filepath} not found.")
            return [{}] 
        
        except Exception as e:
            logging.error(f"Unexpected error reading DOCX file: {e}")
            return [{}]      
        
        # ------------- STEP 2: 파싱 시작 
        records: List[Dict[str, Any]] = []      # 앞으로 쓰일 레코드
        table_info_cache: Dict[str, Any] = {}   # 테이블 저장을 위한 캐시 딕셔너리

        previous_title = "" # 최근에 탐색된 Table을 [table_title] 메타데이터에 에 삽입하기 위함
        
        for el in elements:
            record_id = str(uuid.uuid4())
            
            # -------- 레코드의 형식은 다음과 같음
            rec = {
                "id": record_id,
                "category": el.category,
                "context": el.text,
                "metadata": {
                    "general_info": el.metadata.to_dict() if el.metadata else {},
                    "table_info": None # 초기에는 None 또는 빈 딕셔너리로 설정한 뒤, Category가 Table인 경우 채워짐
                }
            }

            # -------- 조건문 추가: 테이블 카테고리인 경우 테이블 정보 추출
            if rec.get("category") == "Table" and rec["metadata"]["general_info"].get("text_as_html"):
                table_html_code = rec["metadata"]["general_info"]["text_as_html"]
                table_uid = uuid.uuid4().hex[:8] # 테이블 자체의 고유 ID

                # general_info에서 테이블 관련 원본 HTML 정보 삭제 (text_as_html은 table_info로 이동)
                rec["metadata"]["general_info"].pop("text_as_html", None)
                
                # table_info에 필요한 정보 추가 원본 HTML, 고유 ID, 그리고 현재 제목
                rec["metadata"]["table_info"] = {
                    "table_id": table_uid,
                    "text_as_html": table_html_code, # 원본 HTML은 보관 (후속 처리에서 사용)
                    "html_path": None, # HTML 파일 저장 경로는 RecordManager에서 결정
                    "table_title": previous_title,          
                }

                # 테이블 정보 캐시 (아래서 테이블 HTML 파일 저장 및 경로 업데이트)
                table_info_cache[table_uid] = rec["metadata"]["table_info"]

            records.append(rec)
            
            if rec.get("category") == "Title": # 위에 선언한 previous_title 변수에 현재 제목 저장
                previous_title = (rec.get("context") or "").strip()

        logging.info(f"Parsed {len(records)} elements from DOCX: {self.filepath.name}")

        # ------------- STEP 3: 테이블 HTML 파일 저장 및 경로 업데이트 (TODO: 챗봇에서 표를 출력하고 싶을 때 사용, 마크다운이 될 수 있음)
        
        for table_uid, table_info in table_info_cache.items():
                html_code = table_info.get("text_as_html")
                if html_code:
                    html_file_path = self._save_table_html(html_code, table_uid)
                    # 해당 테이블 레코드의 html_path 업데이트
                    for rec in records: # RecordProcessor 내부 레코드 접근
                        if rec.get("category") == "Table" and rec.get("metadata", {}).get("table_info", {}).get("table_id") == table_uid:
                            rec["metadata"]["table_info"]["html_path"] = str(html_file_path)
                            break # 테이블 레코드를 찾으면 루프 종료

        logging.info(f"Table HTML {len(table_info_cache)} elements saved")

        self.records = records  # 클래스 속성에 저장
        self.logger.info(f"Total records parsed: {len(self.records)}")

        return self.records
        
    def preprocess(self) -> List[Dict[str, Any]]:
        # ------------- STEP 1: 이미지 추출 (선택)
        self._save_images()
        
        # ------------- STEP 2: 섹션 정보 메타데이터 주입 (max_depth 깊이까지 추가)
        self.add_sections_to_metadata(max_depth = 3)

        return self.records
    
    def get_chunking_configs(self) -> List[Dict[str, Any]]:
        """ DOCX 파일에 대한 기본 청킹 설정 목록을 반환합니다. """
        return [
            {
                "chunking_type": ChunkingType.NARRATIVE,
                "strategy": ChunkingStrategy.RECURSIVE,
                "kwargs": {"chunk_size": 500, "chunk_overlap": 50}
            },
            {
                "chunking_type": ChunkingType.TABLE,
                "strategy": ChunkingStrategy.ROW,
                "kwargs": {}
            }
            # 필요한 경우 DOCX에 특화된 다른 청킹 설정 추가 가능
            # {
            #     "chunking_type": ChunkingType.NARRATIVE,
            #     "strategy": TextChunkingStrategy.FIXED,
            #     "kwargs": {"chunk_size": 200}
            # }
        ]
    
    def postprocess(self) -> List[Dict[str, Any]]:
        # ------------- STEP 1: 메타데이터 삭제
        self.delete_metadata(["category_depth", "page_number", "filetype", "languages","parent_id"])
        
        # ------------- STEP 2: 타이틀 삭제
        self.remove_title_records()
        
        return self.records
    
    # ----------------------------------
    #         클래스 내부 구현 함수
    # ----------------------------------    
    def add_sections_to_metadata(self, max_depth: int = 5) -> None:
        """
        STEP 2: 섹션 정보 메타데이터 주입
            - category == "Title" & general_info.category_depth(0~5)로 섹션 스택 갱신
            - general_info에 section0~section{max_depth-1}, section_path 추가
            - 누락/잘못된 필드는 전부 None으로 채움
        """
        section_stack: List[Optional[str]] = [None] * max_depth # 섹션 스택 초기화

        for rec in self.records:
            general_info = rec.get("metadata", {}).get("general_info", {})

            if rec.get("category") == "Title":
                try:
                    d = int(general_info.get("category_depth"))
                except (TypeError, ValueError):
                    d = None

                # 카테고리 뎁스가 None이 아니고, max_depth 범위 내에 있을 경우에만 처리
                if d is not None and 0 <= d < max_depth:
                    section_stack[d] = (rec.get("context") or "").strip() or None
                    # d 이후의 인덱스들은 None 으로 설정
                    for i in range(d + 1, max_depth):
                        section_stack[i] = None

            # section_path 생성
            general_info["section_path"] = " ‣ ".join([s for s in section_stack if s]) or None

            # section0~section{max_depth-1} 필드 추가
            for i in range(max_depth):
                general_info[f"section{i}"] = section_stack[i]

            # 업데이트 된 정보 다시 rec에 저장
            rec["metadata"]["general_info"] = general_info

        self.logger.info(f"sections metadata added. max depth: {max_depth}.")

    def delete_metadata(self, elements: List[str]) -> None:
        """
        지정된 메타데이터 키를 오리지널 레코드에서 삭제
            - elements: 삭제할 키 목록 (예: ["category_depth", "section0"])
            - 라이브러리에서 나온 메타데이터기 때문에 general_info와 table_info에서 확인
        """
        for rec in self.records:
            general_info = rec.get("metadata", {}).get("general_info", {})
            table_info = rec.get("metadata", {}).get("table_info", {})
            chunk_info = rec.get("metadata", {}).get("chunk_info", {})

            for key in elements:
                general_info.pop(key, None)
                if table_info:
                    table_info.pop(key, None)

            rec["metadata"] = {
                "general_info": general_info,
                "table_info": table_info if table_info else None,
                "chunk_info": chunk_info,
            }
        
        self.logger.info(f"다음 키를 메타데이터에서 삭제했습니다. '{elements}'")

    def remove_title_records(self) -> None:
        """
        'category'가 'Title'인 레코드들을 self.records에서 제거합니다.
        """
        initial_count = len(self.records)
        self.records = [rec for rec in self.records if rec.get("category") != "Title"]
        removed_count = initial_count - len(self.records)
        self.logger.info(f"'Title' 카테고리 레코드 {removed_count}개를 제거했습니다.")

class PptxPreprocessor(BasePreprocessor):
    def __init__(self, filename: str):
        super().__init__(filename)

    # ----------------------------------
    #           상속받은 메소드
    # ----------------------------------    
    def parse(self) -> List[Dict[str, Any]]:
        """
        PPTX 파일을 파싱하여 DocumentRecord 리스트와 table_info_cache를 반환합니다.
        section_path는 "슬라이드 제목 ‣ 소제목" 형식이며 section0~section4 필드도 함께 저장됩니다.
        """
        try:
            prs = Presentation(str(self.filepath))
        except Exception as e:
            logging.error(f"Error reading PPTX file: {e}")
            return [{}]

        records: List[Dict[str, Any]] = []
        table_info_cache: Dict[str, Any] = {}

        for slide in prs.slides:
            slide_title = None
            sub_title = None

            # 제목과 소제목 감지 (placeholder 우선)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if not shape.is_placeholder:
                    continue

                try:
                    phf_type = shape.placeholder_format.type
                    text = shape.text.strip()
                    if not text:
                        continue

                    if phf_type.name == "TITLE" and not slide_title:
                        slide_title = text
                    elif phf_type.name == "SUBTITLE" and not sub_title:
                        sub_title = text
                except Exception:
                    continue

            # section path 구성
            section_titles = [slide_title, sub_title]
            section_path_str = " › ".join([s for s in section_titles if s])
            section_fields = {
                f"section{i}": section_titles[i] if i < len(section_titles) else None
                for i in range(5)
            }

            # shape 순회: 텍스트, 표 등
            for shape in slide.shapes:
                if not shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                    continue

                record_id = str(uuid.uuid4())

                if shape.has_text_frame:
                    text = shape.text.strip()
                    if not text or text in section_titles:
                        continue  # 제목/소제목은 context로 저장하지 않음
                
                # -------- 레코드의 형식은 다음과 같음
                    rec = {
                        "id": record_id,
                        "category": "Text",
                        "context": text,
                        "metadata": {
                            "general_info": {
                                "section_path": section_path_str if section_path_str else None,
                                **section_fields
                            },
                            "table_info": None
                        }
                    }
                    records.append(rec)
                
                # -------- 조건문 추가: 테이블 카테고리인 경우 테이블 정보 추출
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    html = "<table>\n"
                    for row in table.rows:
                        html += "  <tr>" + "".join(f"<td>{cell.text.strip()}</td>" for cell in row.cells) + "</tr>\n"
                    html += "</table>"

                    table_uid = uuid.uuid4().hex[:8]
                    rec = {
                        "id": record_id,
                        "category": "Table",
                        "context": "",
                        "metadata": {
                            "general_info": {
                                "section_path": section_path_str if section_path_str else None,
                                **section_fields
                            },
                            "table_info": {
                                "table_id": table_uid,
                                "text_as_html": html,
                                "html_path": None,
                                "table_title": sub_title or slide_title or ""
                            }
                        }
                    }
                    records.append(rec)

                    table_info_cache[table_uid] = rec["metadata"]["table_info"]

        logging.info(f"Parsed {len(records)} elements from PPTX: {self.filepath.name}")
        
        # ------------- STEP 3: 테이블 HTML 파일 저장 및 경로 업데이트 (TODO: 챗봇에서 표를 출력하고 싶을 때 사용, 마크다운이 될 수 있음)
        for table_uid, table_info in table_info_cache.items():
                html_code = table_info.get("text_as_html")
                if html_code:
                    html_file_path = self._save_table_html(html_code, table_uid)
                    # 해당 테이블 레코드의 html_path 업데이트
                    for rec in records: # RecordProcessor 내부 레코드 접근
                        if rec.get("category") == "Table" and rec.get("metadata", {}).get("table_info", {}).get("table_id") == table_uid:
                            rec["metadata"]["table_info"]["html_path"] = str(html_file_path)
                            break # 테이블 레코드를 찾으면 루프 종료

        logging.info(f"Table HTML {len(table_info_cache)} elements saved")

        self.records = records  # 클래스 속성에 저장
        self.logger.info(f"Total records parsed: {len(self.records)}")
        
        return self.records

    def get_chunking_configs(self) -> List[Dict[str, Any]]:
        return super().get_chunking_configs()
    def preprocess(self) -> List[Dict[str, Any]]:
        return super().preprocess()
    def postprocess(self) -> List[Dict[str, Any]]:
        return super().postprocess()

class PreprocessingPipeline:
    """
    PreprocessingPipeline은 파일 확장자에 따라 DOCX 또는 PPTX 파일을 전처리하는 파이프라인입니다.
    """
    def __init__(self, filename: str):
        self.logger = logging.getLogger(__name__)
        self.filename = filename
        self.preprocessor: BasePreprocessor # 타입을 BasePreprocessor로 명시하여 어떤 전처리기든 받을 수 있게 함
        self._initialize_preprocessor()

    def _initialize_preprocessor(self):
        """
        파일 확장자에 따라 적절한 Preprocessor 인스턴스를 생성합니다.
        """
        file_extension = os.path.splitext(self.filename)[1].lower()

        if file_extension == '.docx':
            self.preprocessor = DocxPreprocessor(self.filename)
            self.logger.info(f"[{self.__class__.__name__}] DOCX 파일 감지: {self.filename} - DocxPreprocessor를 사용합니다.")

        elif file_extension == '.pptx':
            self.preprocessor = PptxPreprocessor(self.filename)
            self.logger.info(f"[{self.__class__.__name__}] PPTX 파일 감지: {self.filename} - PptxPreprocessor를 사용합니다.")

        else:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {file_extension}. 지원되는 형식은 .docx, .pptx입니다.")

    def run(self) -> None:
        """
        전처리 파이프라인을 실행합니다.
        1. parse: DOCX 파일을 파싱하여 레코드 생성
        2. preprocess: 이미지 추출 및 섹션 정보 메타데이터 주입
        3. chunk: 청킹 작업 수행
        4. postprocess: 메타데이터 삭제 및 타이틀 레코드 제거
        """
        self.logger.info(f"전처리 파이프라인 시작: {self.filename}")
        self.preprocessor.parse()
        self.preprocessor.preprocess()

        chunking_configs = self.preprocessor.get_chunking_configs()
        for i, config in enumerate(chunking_configs):
            chunking_type   = config["chunking_type"]
            strategy        = config["strategy"]
            kwargs          = config.get("kwargs", {})
            
            self.logger.info(f"청킹 단계 {i+1}/{len(chunking_configs)}: 유형='{chunking_type.name}', 전략='{strategy.name}'")
            self.preprocessor.chunk(
                chunking_type   =   chunking_type,
                strategy        =   strategy,
                **kwargs
            )

        self.preprocessor.postprocess()
        self.preprocessor.to_jsonl()
        
    
if __name__ == "__main__":
    preprocessing = PreprocessingPipeline(filename = settings.FILENAME)
    preprocessing.run()
